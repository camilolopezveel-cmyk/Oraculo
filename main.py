import os
import json
import time
import tempfile
import pyttsx3
import webbrowser
import subprocess
import speech_recognition as sr
from faster_whisper import WhisperModel
from groq import Groq
from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt
import datetime
import re
import base64

# ==========================================
# 1. Configuración de API Keys y Modelos
# ==========================================
# Reemplaza 'TU_API_KEY_AQUI' con tu clave real de Groq o usa un archivo .env
from dotenv import load_dotenv
load_dotenv()

API_KEY = os.environ.get("GROQ_API_KEY", "TU_API_KEY_AQUI")
os.environ["GROQ_API_KEY"] = API_KEY

try:
    # Pasamos la clave explícitamente para evitar problemas
    groq_client = Groq(api_key=API_KEY)
except Exception as e:
    print("Error al inicializar Groq. Asegúrate de configurar GROQ_API_KEY.")

# ==========================================
# 2. Configuración de Voz (Text-to-Speech)
# ==========================================
# Usamos pyttsx3 porque funciona offline, consume 0% de internet
# y muy poca CPU (ideal para PC modesta).
engine = pyttsx3.init()
engine.setProperty('rate', 170) # Velocidad del habla
# Opcional: Cambiar voz a español si hay varias instaladas en Windows
# voices = engine.getProperty('voices')
# for voice in voices:
#     if "spanish" in voice.name.lower() or "ES" in voice.id:
#         engine.setProperty('voice', voice.id)
#         break

def speak(text):
    """Reproduce el texto mediante voz y lo imprime."""
    print(f"\nOráculo: {text}")
    engine.say(text)
    engine.runAndWait()

# ==========================================
# 3. Configuración de Oído (Speech-to-Text)
# ==========================================
print("Cargando modelo Whisper 'tiny' en CPU (optimizado para PC modesta)...")
# Usamos int8 en CPU para que el consumo de RAM sea mínimo (~300MB)
whisper_model = WhisperModel("tiny", device="cpu", compute_type="int8")
recognizer = sr.Recognizer()

def listen():
    """Escucha el micrófono y transcribe el audio usando faster-whisper."""
    with sr.Microphone() as source:
        print("\nAjustando ruido de fondo...")
        recognizer.adjust_for_ambient_noise(source, duration=1)
        print(">> Escuchando... (habla ahora)")
        try:
            # Capturar audio del micrófono
            audio = recognizer.listen(source, timeout=5, phrase_time_limit=10)
            
            # Guardar el audio temporalmente para procesarlo
            with tempfile.NamedTemporaryFile(suffix=".wav", delete=False) as tmp_file:
                tmp_file.write(audio.get_wav_data())
                tmp_filename = tmp_file.name
                
            print("⏳ Procesando audio...")
            segments, _ = whisper_model.transcribe(tmp_filename, beam_size=5, language="es")
            
            text = "".join([segment.text for segment in segments])
            os.remove(tmp_filename) # Limpiar archivo temporal
            
            if text.strip():
                print(f">> Tú: {text.strip()}")
            return text.strip()
            
        except sr.WaitTimeoutError:
            return "" # No se escuchó nada
        except Exception as e:
            print(f"Error al escuchar: {e}")
            return ""

# ==========================================
# 4. Personalidad y LLM (Groq)
# ==========================================
# System Prompt para una personalidad 'sarcástica y eficiente'
SYSTEM_PROMPT = """Eres 'Oráculo', un profesor emérito y experto universal en cualquier disciplina o materia académica sobre la que te consulten. 
Tu trato es amable y respetuoso, enfocado estrictamente en lo académico. Tu misión es que el estudiante saque la máxima calificación.
Como experto:
1. Usas rigor académico y científico, pero explicas con gran claridad pedagógica.
2. Si el usuario pide un DOCUMENTO: Responde con Markdown detallado, con información precisa y estructura académica de la UPANA.
3. Si el usuario pide una PRESENTACIÓN: Responde con una lista estructurada de diapositivas usando este formato exacto:
   DIAPOSITIVA 1: [Título] | [Contenido]
   DIAPOSITIVA 2: [Título] | [Contenido]
   ... y así sucesivamente.
Mantén un tono profesional, claro y directo, evitando el exceso de motivación o frases de ánimo innecesarias."""

# Memoria de la conversación
conversation_history = [
    {"role": "system", "content": SYSTEM_PROMPT}
]

TOOLS = [
    {
        "type": "function",
        "function": {
            "name": "run_powershell",
            "description": "Ejecuta un comando en PowerShell de Windows. Úsalo SIEMPRE que necesites interactuar con la computadora local.",
            "parameters": {
                "type": "object",
                "properties": {
                    "command": {
                        "type": "string",
                        "description": "El código exacto de PowerShell a ejecutar."
                    }
                },
                "required": ["command"]
            }
        }
    }
]

def encode_image(image_path):
    with open(image_path, "rb") as image_file:
        return base64.b64encode(image_file.read()).decode('utf-8')

def extract_image_paths(text):
    pattern1 = r'"([A-Za-z]:[\\/][^"<>|?*]+\.(?:jpg|jpeg|png|webp|gif|bmp))"'
    pattern2 = r'([A-Za-z]:[\\/][^\s"<>|?*]+\.(?:jpg|jpeg|png|webp|gif|bmp))'
    paths = re.findall(pattern1, text, re.IGNORECASE) + re.findall(pattern2, text, re.IGNORECASE)
    valid_paths = []
    for p in paths:
        if os.path.isfile(p):
            valid_paths.append(p)
    return list(set(valid_paths))

def think(user_text):
    """Procesa el texto con el modelo Llama 3 en Groq y devuelve la respuesta."""
    image_paths = extract_image_paths(user_text)
    
    if image_paths:
        img_path = image_paths[0]
        try:
            base64_image = encode_image(img_path)
            
            ext = img_path.split('.')[-1].lower()
            mime_type = "image/jpeg"
            if ext == "png": mime_type = "image/png"
            elif ext == "webp": mime_type = "image/webp"
            elif ext == "gif": mime_type = "image/gif"
            
            vision_history = list(conversation_history)
            vision_history.append({
                "role": "user",
                "content": [
                    {"type": "text", "text": user_text},
                    {
                        "type": "image_url",
                        "image_url": {
                            "url": f"data:{mime_type};base64,{base64_image}",
                        },
                    },
                ]
            })
            
            response = groq_client.chat.completions.create(
                model="meta-llama/llama-4-scout-17b-16e-instruct",
                messages=vision_history,
                temperature=0.7,
                max_tokens=8000
            )
            reply = response.choices[0].message.content
            conversation_history.append({"role": "user", "content": f"{user_text} [Imagen analizada: {img_path}]"})
            conversation_history.append({"role": "assistant", "content": reply})
            return reply
            
        except Exception as e:
            print(f"Error procesando la imagen: {e}")
            return "Intenté analizar la imagen, pero hubo un error. ¿Seguro que la ruta es correcta y la imagen no es demasiado grande?"

    conversation_history.append({"role": "user", "content": user_text})
    
    while True: # Bucle para procesar llamadas a herramientas
        try:
            response = groq_client.chat.completions.create(
                model="llama-3.3-70b-versatile", # Modelo tope de gama de Meta
                messages=conversation_history,
                temperature=0.7,
                max_tokens=8000,
                tools=TOOLS,
                tool_choice="auto"
            )
            
            message = response.choices[0].message
            
            # Si el modelo decidió no usar herramientas, simplemente responde
            if not message.tool_calls:
                reply = message.content
                conversation_history.append({"role": "assistant", "content": reply})
                return reply
                
            # Si decidió usar herramientas, procesamos cada una
            conversation_history.append(message) # Guardar el intento
            
            for tool_call in message.tool_calls:
                if tool_call.function.name == "run_powershell":
                    args = json.loads(tool_call.function.arguments)
                    command = args.get("command", "")
                    
                    print(f"\n[⚠️ ALERTA DE SEGURIDAD] Oráculo quiere ejecutar en tu PC:")
                    print(f"Comando: {command}")
                    confirm = input("¿Permitir ejecución? (S/N): ")
                    
                    if confirm.lower() == 's':
                        print("Ejecutando...")
                        try:
                            # Ejecutar comando en PowerShell
                            result = subprocess.run(
                                ["powershell", "-Command", command], 
                                capture_output=True, text=True, timeout=30
                            )
                            output = result.stdout if result.stdout else result.stderr
                            if not output.strip():
                                output = "Comando ejecutado con éxito sin salida de consola."
                        except Exception as ex:
                            output = f"Error de sistema: {ex}"
                    else:
                        output = "ERROR: El usuario denegó el permiso por seguridad. No pudiste realizar la tarea."
                        print("Operación cancelada.")
                        
                    # Enviar el resultado de vuelta al cerebro
                    conversation_history.append({
                        "role": "tool",
                        "tool_call_id": tool_call.id,
                        "name": tool_call.function.name,
                        "content": str(output)
                    })
                    
        except Exception as e:
            print(f"Error en el cerebro de Groq: {e}")
            return "Oh, parece que tuve un pequeño lapsus en la conexión. ¿Podemos intentarlo de nuevo, por favor?"

# ==========================================
# 5. Ejecución de Comandos Básicos
# ==========================================
def create_word_document(topic):
    speak(f"Redactando un documento sobre {topic}. Dale un momento a mi cerebro...")
    
    prompt = f"Escribe un ensayo detallado y profesional sobre el siguiente tema: {topic}. No incluyas saludos ni comentarios adicionales, solo el contenido del texto."
    try:
        # Llamamos a Groq para generar el texto extenso
        response = groq_client.chat.completions.create(
            model="llama-3.3-70b-versatile",
            messages=[{"role": "user", "content": prompt}],
            temperature=0.7,
            max_tokens=8000
        )
        content = response.choices[0].message.content
        
        # Crear el documento Word
        doc = Document()
        doc.add_heading(topic.title(), 0)
        doc.add_paragraph(content)
        
        # Guardar en el escritorio
        filename = f"Documento_{datetime.datetime.now().strftime('%H%M%S')}.docx"
        filepath = os.path.join(os.path.expanduser("~"), "Desktop", filename)
        doc.save(filepath)
        
        speak(f"Trabajo terminado. He guardado el archivo en tu escritorio como {filename}.")
    except Exception as e:
        print(f"Error al crear documento: {e}")
        speak("Tuve un pequeño problema técnico al generar el documento, pero no te preocupes, lo podemos intentar de nuevo.")

def create_powerpoint_presentation(topic):
    speak(f"Preparando una presentación sobre {topic}. Por favor espera...")
    
    prompt = f"Escribe el contenido para una presentación de diapositivas sobre: {topic}. Formato estricto: Para cada diapositiva empieza con 'DIAPOSITIVA: [Título de la diapositiva]' seguido de una nueva línea y luego los puntos principales (con guiones). No incluyas introducciones ni despedidas, solo el texto de las diapositivas."
    try:
        response = groq_client.chat.completions.create(
            model="llama-3.3-70b-versatile",
            messages=[{"role": "user", "content": prompt}],
            temperature=0.7,
            max_tokens=8000
        )
        content = response.choices[0].message.content
        
        prs = Presentation()
        slides_data = content.split('DIAPOSITIVA:')
        for slide_text in slides_data:
            if not slide_text.strip():
                continue
                
            lines = [line.strip() for line in slide_text.strip().split('\n') if line.strip()]
            if not lines: continue
            
            title = lines[0].strip('[]')
            body_points = lines[1:]
            
            slide_layout = prs.slide_layouts[1] # Title and Content
            slide = prs.slides.add_slide(slide_layout)
            
            slide.shapes.title.text = title
            tf = slide.placeholders[1].text_frame
            
            for i, point in enumerate(body_points):
                clean_point = point.lstrip('-*• ')
                if i == 0:
                    tf.text = clean_point
                else:
                    p = tf.add_paragraph()
                    p.text = clean_point
                    
        filename = f"Presentacion_{datetime.datetime.now().strftime('%H%M%S')}.pptx"
        filepath = os.path.join(os.path.expanduser("~"), "Desktop", filename)
        prs.save(filepath)
        
        speak(f"Trabajo terminado. He guardado la presentación en tu escritorio como {filename}.")
    except Exception as e:
        print(f"Error al crear presentación: {e}")
        speak("Tuve un pequeño problema técnico al generar la presentación, pero lo podemos intentar de nuevo.")

def execute_command(text):
    """Busca comandos locales en el texto del usuario para ejecutarlos."""
    text_lower = text.lower()
    
    # 1. Abrir Navegador
    if "abrir navegador" in text_lower or "abre google" in text_lower:
        speak("Abriendo tu navegador para la investigación.")
        webbrowser.open("https://www.google.com")
        return True
        
    # 2. Buscar en YouTube
    elif "buscar en youtube" in text_lower:
        query = text_lower.replace("buscar en youtube", "").strip()
        if query:
            speak(f"Buscando '{query}' en YouTube.")
            webbrowser.open(f"https://www.youtube.com/results?search_query={query}")
        else:
            speak("Abriendo YouTube. ¿En qué te puedo ayudar a buscar hoy?")
            webbrowser.open("https://www.youtube.com")
        return True
        
    # 3. Abrir aplicación local (Ejemplo: Bloc de notas)
    elif "abrir bloc de notas" in text_lower:
        speak("Abriendo el bloc de notas para registrar información.")
        subprocess.Popen(["notepad.exe"])
        return True
        
    # 4. Crear documento Word
    elif any(phrase in text_lower for phrase in ["crear documento", "escribe un documento", "crea un documento", "crear un documento", "haz un documento", "hacer un documento"]):
        # Extraer el tema de la frase
        topic = text_lower
        for phrase in ["crear documento sobre", "crear documento de", "crear documento", "escribe un documento sobre", "escribe un documento de", "escribe un documento", "crea un documento sobre", "crea un documento de", "crea un documento", "crear un documento sobre", "crear un documento de", "crear un documento", "haz un documento sobre", "haz un documento de", "haz un documento", "hacer un documento sobre", "hacer un documento de", "hacer un documento"]:
            topic = topic.replace(phrase, "")
        topic = topic.strip()
        if not topic:
            topic = "Inteligencia Artificial" # Tema por defecto si no especifica
        create_word_document(topic)
        return True
        
    # 5. Crear presentacion
    elif any(phrase in text_lower for phrase in ["crear presentacion", "crear presentación", "haz una presentacion", "haz una presentación", "hacer presentacion", "hacer presentación", "crea una presentacion", "crea una presentación"]):
        topic = text_lower
        for phrase in ["crear presentacion sobre", "crear presentación sobre", "crear presentacion de", "crear presentación de", "haz una presentacion sobre", "haz una presentación sobre", "crea una presentacion sobre", "crea una presentación sobre", "hacer presentacion sobre", "hacer presentación sobre", "crear presentacion", "crear presentación", "haz una presentacion", "haz una presentación", "crea una presentacion", "crea una presentación", "hacer presentacion", "hacer presentación"]:
            topic = topic.replace(phrase, "")
        topic = topic.strip()
        if not topic:
            topic = "Inteligencia Artificial" # Tema por defecto
        create_powerpoint_presentation(topic)
        return True
        
    # 6. Salir
    elif "apagar sistema" in text_lower or "adiós oráculo" in text_lower:
        speak("Ha sido un placer ayudarte. Apagando sistemas. Hasta pronto.")
        return "EXIT"
        
    return False

# ==========================================
# 6. Bucle Principal (Estructura)
# ==========================================
def main():
    if "TU_API_KEY_AQUI" in os.environ.get("GROQ_API_KEY", ""):
        print("\n[!] ADVERTENCIA: No has configurado tu GROQ_API_KEY.")
        print("[!] Modifica el archivo main.py en la línea 14 con tu clave antes de usar el cerebro.\n")
        
    speak("Sistemas inicializados. Oráculo en línea. Hola, ¿en qué te puedo ayudar hoy?")
    
    while True:
        # Modo chat: El usuario escribe su petición
        user_input = input("\nTú: ")
        
        if not user_input.strip():
            continue # Si presionó Enter sin escribir nada, vuelve a preguntar
            
        # 2. Acciones: Comprobar si es un comando del sistema
        command_result = execute_command(user_input)
        
        if command_result == "EXIT":
            break # Termina el programa
        elif command_result == True:
            continue # Si era un comando local, ya se ejecutó, vuelve a escuchar
            
        # 3. Cerebro: Si no es comando, enviar al LLM
        response = think(user_input)
        
        # 4. Voz: Responder al usuario
        speak(response)

if __name__ == "__main__":
    main()
